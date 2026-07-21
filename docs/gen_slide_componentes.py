#!/usr/bin/env python3
"""Genera una slide (PPTX) con los componentes a implementar y para qué se usa cada uno."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

PRIM = RGBColor(0x1D, 0x4E, 0xD8)
DARK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CARDS = [
    ("1 · Motor de cálculo — cierre", RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0xEF, 0xF4, 0xFF), [
        ("Re-export de fórmulas (acumuladores)", "Completar el catálogo real (89 acumuladores truncados) → tasas con importes reales."),
        ("Mapeo de variables @I_* por inmueble", "Derivar zona, frente, esquina, unidades… de los datos del contribuyente, sin cargarlas a mano."),
        ("Confirmar #REDONDEOESPECIAL", "Fijar la regla de redondeo exacta del recibo."),
    ]),
    ("2 · Ciclo de cobranza", RGBColor(0x05, 0x96, 0x69), RGBColor(0xEC, 0xFD, 0xF5), [
        ("Cobro masivo por lote", "Integrar Tesorería con la deuda: cobrar recibos en masa y bajar la cuenta corriente."),
        ("Cuenta corriente / libro mayor", "Debe/Haber completo, saldos a favor y compensación."),
    ]),
    ("3 · Módulos de gestión (nuevos)", RGBColor(0xB4, 0x53, 0x09), RGBColor(0xFF, 0xF7, 0xED), [
        ("Gestión de cobranzas", "Operativos, intimaciones y seguimiento de morosos."),
        ("Apremios / Legales", "Juicio de apremio, honorarios y embargos sobre deuda vencida."),
        ("Débito automático", "Archivos de bancos/tarjetas, presentación y rechazos."),
        ("Informes / Reportes", "Reportería de gestión y exportables (PDF/Excel)."),
    ]),
    ("4 · Frontend / Portal", RGBColor(0x6D, 0x28, 0xD9), RGBColor(0xF5, 0xF3, 0xFF), [
        ("Dashboard con KPIs", "Deuda total, recaudación del día/mes, emisiones pendientes."),
        ("Portal del contribuyente", "Consulta web, boletas y pasarela de pago."),
        ("Importación de lotes", "Cargar y seguir lotes de novedades/pagos."),
    ]),
]


def _no_line(shape):
    shape.line.fill.background()


def add_card(slide, x, y, w, h, titulo, color, fill_bg, items):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = fill_bg
    card.line.color.rgb = color; card.line.width = Pt(1)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_right = Pt(12); tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    r = p.add_run(); r.text = titulo
    r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = color
    p.space_after = Pt(4)

    for comp, uso in items:
        pc = tf.add_paragraph()
        rc = pc.add_run(); rc.text = "• " + comp
        rc.font.bold = True; rc.font.size = Pt(10.5); rc.font.color.rgb = DARK
        pc.space_after = Pt(0)
        pu = tf.add_paragraph()
        ru = pu.add_run(); ru.text = "   " + uso
        ru.font.size = Pt(9); ru.font.color.rgb = GREY
        pu.space_after = Pt(5)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Banda de título
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
band.fill.solid(); band.fill.fore_color.rgb = DARK; _no_line(band); band.shadow.inherit = False
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.14), Inches(12.3), Inches(0.95)).text_frame
tb.word_wrap = True
p = tb.paragraphs[0]
r = p.add_run(); r.text = "Cheyenne · Componentes a implementar"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tb.add_paragraph()
r2 = p2.add_run(); r2.text = "Qué falta construir y para qué se usa cada uno"
r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xB8, 0xC2, 0xD9)

# Grilla 2x2 de tarjetas
gap = Inches(0.35)
top = Inches(1.45)
cw = (prs.slide_width - Inches(1.0) - gap) / 2
# alturas: fila de arriba más baja (menos items), abajo más alta
rows_h = [Inches(2.45), Inches(3.05)]
xs = [Inches(0.5), Inches(0.5) + cw + gap]
ys = [top, top + rows_h[0] + Inches(0.2)]
pos = [(0, 0), (1, 0), (0, 1), (1, 1)]
for (col, row), card in zip(pos, CARDS):
    add_card(slide, xs[col], ys[row], cw, rows_h[row], card[0], card[1], card[2], card[3])

import os
out = os.path.join(os.path.dirname(__file__), "Cheyenne-Componentes.pptx")
prs.save(out)
print("OK ->", out)
